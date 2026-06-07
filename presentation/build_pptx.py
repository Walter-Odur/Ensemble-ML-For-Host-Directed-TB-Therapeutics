"""
Build a 16-slide scientific presentation for the TB ML manuscript.
Font: Calibri | Background: White | Titles: 18pt | Subtitles: 16pt | Body: 14pt
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

FIGS = os.path.join(os.path.dirname(__file__), '..', 'results', 'figures', 'manuscript')

# Colours
DARK = RGBColor(0x2C, 0x3E, 0x50)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BODY = RGBColor(0x33, 0x33, 0x33)
GREY = RGBColor(0x55, 0x55, 0x55)
ACCENT = RGBColor(0x2C, 0x3E, 0x50)
RED = RGBColor(0xC0, 0x39, 0x2B)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
MUTED = RGBColor(0x7F, 0x8C, 0x8D)

FONT = 'Calibri'

def add_blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def add_header_bar(slide, title_text):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.7))
    shp.fill.solid()
    shp.fill.fore_color.rgb = DARK
    shp.line.fill.background()
    tf = shp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.1)

def add_textbox(slide, left, top, width, height, text, size=14, bold=False, color=BODY, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT
    p.font.italic = italic
    p.alignment = align
    p.space_after = Pt(4)
    return tf

def add_bullet_list(slide, left, top, width, height, items, size=14, color=BODY, bold_prefix=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        p.space_before = Pt(4)
        p.level = 0
        if bold_prefix and ':' in item:
            bold_part, rest = item.split(':', 1)
            run_b = p.add_run()
            run_b.text = '\u2022 ' + bold_part + ':'
            run_b.font.size = Pt(size)
            run_b.font.bold = True
            run_b.font.color.rgb = color
            run_b.font.name = FONT
            run_r = p.add_run()
            run_r.text = rest
            run_r.font.size = Pt(size)
            run_r.font.color.rgb = color
            run_r.font.name = FONT
        else:
            run = p.add_run()
            run.text = '\u2022 ' + item
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.name = FONT
    return tf

def add_image_scaled(slide, img_path, max_w, max_h, left, top):
    from PIL import Image
    im = Image.open(img_path)
    w_px, h_px = im.size
    ar = w_px / h_px
    if max_w / max_h > ar:
        h = max_h; w = h * ar
    else:
        w = max_w; h = w / ar
    slide.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(w), Inches(h))

def add_line(slide, left, top, width):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(2))
    shp.fill.solid()
    shp.fill.fore_color.rgb = DARK
    shp.line.fill.background()

# ============================================================
# SLIDE 1: TITLE
# ============================================================
s = add_blank_slide()
add_textbox(s, 1.5, 1.2, 10.3, 1.8,
    'Single-Cell Transcriptomics of TB Granulomas Reveals\nHost-Directed Drug Candidates Through Ensemble Machine Learning',
    size=22, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_line(s, 5.5, 3.2, 2.3)
add_textbox(s, 1.5, 3.5, 10.3, 0.5, 'WALTER ODUR', size=16, bold=True, color=BODY, align=PP_ALIGN.CENTER)
add_textbox(s, 1.5, 4.1, 10.3, 0.4, '2025/HD07/26017U', size=14, color=GREY, align=PP_ALIGN.CENTER)
add_textbox(s, 1.5, 4.8, 10.3, 0.5,
    'Department of Immunology and Molecular Biology, School of Biomedical Sciences,\nCollege of Health Sciences, Makerere University',
    size=12, color=GREY, align=PP_ALIGN.CENTER, italic=True)
add_textbox(s, 1.5, 5.5, 10.3, 0.5,
    'The African Center of Excellence in Bioinformatics and Data Intensive Sciences,\nMakerere University, Kampala, Uganda',
    size=12, color=GREY, align=PP_ALIGN.CENTER, italic=True)

# ============================================================
# SLIDE 2: BACKGROUND
# ============================================================
s = add_blank_slide()
add_header_bar(s, 'Background')
tf = add_textbox(s, 0.5, 1.0, 12.3, 1.2,
    'Tuberculosis (TB), caused by Mycobacterium tuberculosis, remains the leading infectious disease killer globally, '
    'with 10.6 million new cases and 1.3 million deaths in 2022 (WHO, 2023). Multidrug-resistant strains now account '
    'for nearly 500,000 incident cases annually, and treatment outcomes remain poor.',
    size=14)
tf = add_textbox(s, 0.5, 2.5, 12.3, 1.2,
    'The lung granuloma, a structured aggregate of macrophages, lymphocytes, and fibroblasts, is the hallmark '
    'pathological structure of TB. While it contains infection, it simultaneously provides a niche in which '
    'M. tuberculosis persists in a drug-tolerant state (Ramakrishnan, 2012).',
    size=14)
tf = add_textbox(s, 0.5, 4.0, 12.3, 1.5,
    'Single-cell RNA sequencing (scRNA-seq) has enabled profiling of individual cellular transcriptional states within '
    'infected tissues. However, conventional differential expression analyses fail to capture the nonlinear, '
    'combinatorial patterns in gene expression that distinguish disease from control states. Ensemble machine learning '
    'offers a principled framework for modelling these complex transcriptomic patterns.',
    size=14)

# ============================================================
# SLIDE 3: PROBLEM STATEMENT
# ============================================================
s = add_blank_slide()
add_header_bar(s, 'Problem Statement')
add_bullet_list(s, 0.5, 1.2, 12.3, 5.5, [
    'Drug resistance crisis: Standard TB treatment requires 6 months of antibiotics, and multidrug-resistant strains are increasing, necessitating alternative host-directed therapeutic strategies.',
    'Limited molecular resolution: The granuloma microenvironment harbours transcriptionally diverse cell states that cannot be resolved by bulk analyses, yet single-cell datasets remain underexploited for drug discovery.',
    'Analytical gap: Conventional differential expression methods treat genes independently and miss the nonlinear, combinatorial transcriptomic patterns that ensemble machine learning can capture for therapeutic target identification.'
], size=14)

# ============================================================
# SLIDE 4: OBJECTIVES
# ============================================================
s = add_blank_slide()
add_header_bar(s, 'Study Objectives')
add_textbox(s, 0.5, 1.5, 12.3, 0.4, 'Objective 1', size=16, bold=True, color=ACCENT)
add_textbox(s, 0.5, 2.0, 12.3, 1.0,
    'To develop and benchmark an ensemble machine learning pipeline for classifying TB-affected versus control '
    'cellular states from single-cell transcriptomic data and to extract a SHAP-derived disease gene signature.',
    size=14)
add_textbox(s, 0.5, 3.5, 12.3, 0.4, 'Objective 2', size=16, bold=True, color=ACCENT)
add_textbox(s, 0.5, 4.0, 12.3, 1.0,
    'To leverage the disease signature for computational drug repurposing through connectivity mapping, '
    'nominating candidate host-directed therapies for tuberculosis.',
    size=14)

# ============================================================
# SLIDE 5: METHODOLOGY SCHEMATIC
# ============================================================
s = add_blank_slide()
add_header_bar(s, 'Methodology')
steps_row1 = [
    ('Data Acquisition', 'Seq-Well scRNA-seq\nSCP3227\nn = 19,044 cells'),
    ('Preprocessing', 'QC filtering\nNormalization\n2,497 HVGs'),
    ('Feature Engineering', 'Binary labelling\n80/20 stratified split\nPCA, t-SNE'),
    ('ML Classification', 'RF, XGBoost,\nLightGBM,\nStacking Ensemble'),
]
steps_row2 = [
    ('HP Tuning', 'Bayesian optimization\nOptuna TPE\n50 trials, 5-fold CV'),
    ('Interpretability', 'SHAP TreeExplainer\n100-gene disease\nsignature'),
    ('Drug Repurposing', 'L1000 CMap query\nEnrichment analysis\nHost-directed targets'),
    ('Validation', 'Expression heatmap\nFunctional annotation\nLiterature support'),
]
box_w, box_h = 2.6, 1.5
gap = 0.35
start_x = 0.5
for row_idx, steps in enumerate([steps_row1, steps_row2]):
    y = 1.2 + row_idx * 2.3
    for i, (title, desc) in enumerate(steps):
        x = start_x + i * (box_w + gap)
        shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(box_w), Inches(box_h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
        shp.line.color.rgb = RGBColor(0xBD, 0xC3, 0xC7)
        shp.line.width = Pt(1)
        tf = shp.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = DARK
        run.font.name = FONT
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = desc
        run2.font.size = Pt(10)
        run2.font.color.rgb = GREY
        run2.font.name = FONT
        p2.alignment = PP_ALIGN.CENTER
        # Arrow between boxes (except last)
        if i < 3:
            arrow_x = x + box_w
            arrow_y = y + box_h / 2 - 0.1
            add_textbox(s, arrow_x, arrow_y, gap, 0.3, '\u2192', size=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 6: FIGURE 1 - DATA OVERVIEW
# ============================================================
s = add_blank_slide()
add_header_bar(s, 'Results: Dataset Overview and Quality Control')
fig1_path = os.path.join(FIGS, 'figure1_data_overview.png')
add_image_scaled(s, fig1_path, 7.5, 4.8, 0.3, 1.0)
tf = add_textbox(s, 8.2, 1.0, 4.8, 1.0, 'Figure 1.', size=12, bold=True, color=DARK)
add_textbox(s, 8.2, 1.4, 4.8, 1.0,
    'The dataset comprises 19,044 cells from TB-positive and TB-negative lung tissue.', size=12)
add_textbox(s, 8.2, 2.3, 4.8, 1.0,
    'Class imbalance: 88.6% TB-affected (n = 16,874) vs 11.4% Control (n = 2,170). HIV-only cells (n = 588) excluded.', size=12)
add_textbox(s, 8.2, 3.5, 4.8, 1.5,
    'QC metrics (genes/cell, UMI counts, mitochondrial fraction) showed comparable distributions between conditions, '
    'confirming biological rather than technical signal.', size=12)

# ============================================================
# SLIDE 7: FIGURE 2 - DIMENSIONALITY REDUCTION
# ============================================================
s = add_blank_slide()
add_header_bar(s, 'Results: Dimensionality Reduction')
fig2_path = os.path.join(FIGS, 'figure2_dimensionality_reduction.png')
add_image_scaled(s, fig2_path, 7.8, 3.5, 0.3, 1.0)
tf = add_textbox(s, 8.5, 1.0, 4.5, 0.4, 'Figure 2.', size=12, bold=True, color=DARK)
add_textbox(s, 8.5, 1.4, 4.5, 1.0,
    'PCA on 2,497 HVGs revealed partial but incomplete class separation. PC1 and PC2 captured 8.5% and 5.3% of variance.', size=12)
add_textbox(s, 8.5, 2.5, 4.5, 1.0,
    'Twenty PCs accounted for approximately 30% of total variance with a characteristic elbow at PC3.', size=12)
add_textbox(s, 8.5, 3.4, 4.5, 1.2,
    't-SNE of a 5,000-cell subsample confirmed distinct clusters with partial overlap, indicating that high-dimensional '
    'modelling is required to fully resolve the transcriptomic distinction.', size=12)

# ============================================================
# SLIDE 8: FIGURE 3 - MODEL PERFORMANCE
# ============================================================
s = add_blank_slide()
add_header_bar(s, 'Results: Model Performance')
fig3_path = os.path.join(FIGS, 'figure3_model_performance.png')
add_image_scaled(s, fig3_path, 7.0, 5.5, 0.2, 0.9)
tf = add_textbox(s, 8.0, 1.0, 5.0, 0.4, 'Figure 3.', size=12, bold=True, color=DARK)
add_textbox(s, 8.0, 1.4, 5.0, 1.0,
    'LightGBM achieved the highest discrimination (AUC-ROC = 0.970, F1 = 0.969).', size=12)
add_textbox(s, 8.0, 2.3, 5.0, 1.0,
    'Stacking ensemble performed comparably (AUC = 0.968) but did not surpass LightGBM.', size=12)
add_textbox(s, 8.0, 3.1, 5.0, 1.2,
    'Confusion matrices revealed complementary error profiles: RF had highest recall (0.991); stacking had highest precision (0.987); LightGBM achieved the best balance.', size=12)
add_textbox(s, 8.0, 4.3, 5.0, 1.0,
    '5-fold CV estimates closely matched test performance, indicating minimal overfitting.', size=12)

# ============================================================
# SLIDE 9: HYPERPARAMETER TUNING TABLE
# ============================================================
s = add_blank_slide()
add_header_bar(s, 'Results: Hyperparameter Tuning')
add_textbox(s, 0.5, 1.0, 12.3, 0.8,
    'Bayesian optimization (Optuna, 50 trials, 5-fold CV) was applied to all four models. '
    'The tuned configurations did not improve over baseline defaults.',
    size=14)

from pptx.util import Inches, Pt
from pptx.oxml.ns import qn

rows, cols = 5, 4
tbl_w = Inches(9.0)
tbl_h = Inches(2.2)
left = Inches(2.1)
top = Inches(2.2)
table_shape = s.shapes.add_table(rows, cols, left, top, tbl_w, tbl_h)
tbl = table_shape.table
col_widths = [Inches(2.8), Inches(2.1), Inches(2.1), Inches(2.0)]
for i, w in enumerate(col_widths):
    tbl.columns[i].width = w

headers = ['Model', 'Baseline AUC', 'Tuned AUC', '\u0394AUC']
data = [
    ['Random Forest', '0.944', '0.926', '\u20130.019'],
    ['XGBoost', '0.962', '0.961', '\u20130.001'],
    ['LightGBM', '0.970', '0.967', '\u20130.003'],
    ['Stacking Ensemble', '0.968', '0.964', '\u20130.004'],
]
for j, h in enumerate(headers):
    cell = tbl.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = FONT
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK

delta_colors = [RED, MUTED, ORANGE, ORANGE]
for i, row_data in enumerate(data):
    for j, val in enumerate(row_data):
        cell = tbl.cell(i + 1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.name = FONT
            p.alignment = PP_ALIGN.CENTER
            if j == 3:
                p.font.color.rgb = delta_colors[i]
                p.font.bold = True

add_textbox(s, 0.5, 4.8, 12.3, 1.5,
    'Default hyperparameter configurations are already well suited to the structure of single-cell transcriptomic data. '
    'The absence of improvement argues against overfitting in the baseline models. '
    'Baseline configurations were retained for all subsequent analyses.',
    size=14)

# ============================================================
# SLIDE 10: FIGURE 4 - SHAP DISEASE SIGNATURE
# ============================================================
s = add_blank_slide()
add_header_bar(s, 'Results: SHAP Disease Signature')
fig4_path = os.path.join(FIGS, 'figure4_disease_signature.png')
add_image_scaled(s, fig4_path, 7.0, 5.5, 0.2, 0.9)
tf = add_textbox(s, 8.0, 1.0, 5.0, 0.4, 'Figure 4.', size=12, bold=True, color=DARK)
add_textbox(s, 8.0, 1.4, 5.0, 1.0,
    'SHAP analysis of LightGBM identified a 100-gene TB disease signature (59 up, 41 down).', size=12)
add_textbox(s, 8.0, 2.3, 5.0, 1.5,
    'Top genes: MTRNR2L12 (mitochondrial stress), TAOK1 (MAPK signalling), S100A12 (antimicrobial calgranulin), '
    'CCL18 (macrophage activation).', size=12)
add_textbox(s, 8.0, 3.8, 5.0, 1.5,
    'Downregulated: RGS2 (chemokine signalling), PLTP and APOC1 (lipid metabolism), '
    'reflecting metabolic reprogramming of granuloma macrophages.', size=12)

# ============================================================
# SLIDE 11: FIGURE 6 - EXPRESSION HEATMAP
# ============================================================
s = add_blank_slide()
add_header_bar(s, 'Results: Expression Validation')
fig6_path = os.path.join(FIGS, 'figure6_expression_heatmap.png')
add_image_scaled(s, fig6_path, 7.5, 5.0, 0.2, 0.9)
tf = add_textbox(s, 8.2, 1.0, 4.8, 0.4, 'Figure 5.', size=12, bold=True, color=DARK)
add_textbox(s, 8.2, 1.4, 4.8, 1.2,
    'Z-scored expression heatmap for the top 30 SHAP genes across 200 Control and 200 TB-affected cells.', size=12)
add_textbox(s, 8.2, 2.6, 4.8, 1.5,
    'Upregulated genes (TAOK1, S100A12, CCL18) showed elevated expression in TB-affected cells. '
    'Downregulated genes (RGS2, SPARC) were higher in controls.', size=12)
add_textbox(s, 8.2, 4.1, 4.8, 1.2,
    'Cell-to-cell variability within TB-affected cells reflects the heterogeneity of the granuloma microenvironment '
    '(Gideon et al., 2022).', size=12)

# ============================================================
# SLIDE 12: FIGURE 5 - DRUG REPURPOSING
# ============================================================
s = add_blank_slide()
add_header_bar(s, 'Results: Drug Repurposing')
fig5_path = os.path.join(FIGS, 'figure5_drug_repurposing.png')
add_image_scaled(s, fig5_path, 7.2, 5.0, 0.2, 0.9)
tf = add_textbox(s, 8.0, 1.0, 5.0, 0.4, 'Figure 6.', size=12, bold=True, color=DARK)
add_textbox(s, 8.0, 1.4, 5.0, 1.2,
    'L1000 Connectivity Map query identified 1,831 compounds at adjusted P < 0.05, with 208 achieving P < 0.0001.', size=12)
add_textbox(s, 8.0, 2.6, 5.0, 1.5,
    'Top candidates: adenosine triphosphate (purinergic signalling), monensin (antimycobacterial ionophore), '
    'cephaeline and anisomycin (protein synthesis inhibitors).', size=12)
add_textbox(s, 8.0, 4.1, 5.0, 1.5,
    'Several candidates have independent preclinical evidence: imatinib (macrophage activation), '
    'metformin (autophagy), verapamil (efflux pump inhibition).', size=12)

# ============================================================
# SLIDE 13: DISCUSSION
# ============================================================
s = add_blank_slide()
add_header_bar(s, 'Discussion')
add_bullet_list(s, 0.5, 1.0, 12.3, 6.0, [
    'High classification accuracy: LightGBM discriminated TB-affected from control cells with AUC-ROC of 0.970, demonstrating that a single well-configured model can approach ceiling performance on scRNA-seq data.',
    'Biologically coherent signature: The SHAP-derived 100-gene signature recapitulates known TB immunopathology, including calgranulin-mediated defence (S100A12), alternative macrophage activation (CCL18), and granuloma fibrotic remodelling (COL3A1, DCN).',
    'Robust defaults: Hyperparameter tuning did not improve performance, indicating that default ML library configurations are already well suited to single-cell transcriptomic data structure.',
    'Pharmacological plausibility: Top drug candidates (monensin, imatinib, metformin, verapamil) have independent preclinical evidence of anti-TB activity, validating the connectivity mapping approach.',
    'Limitations: The classification distinguishes cells from TB-positive vs TB-negative donors rather than infected vs bystander cells. The Cancer Control reference group may carry its own transcriptomic alterations. Drug candidates require experimental validation.'
], size=13)

# ============================================================
# SLIDE 14: FUTURE DIRECTIONS & CONCLUSIONS
# ============================================================
s = add_blank_slide()
add_header_bar(s, 'Future Directions and Conclusions')
add_textbox(s, 0.5, 1.1, 12.3, 0.4, 'Future Directions', size=16, bold=True, color=ACCENT)
add_bullet_list(s, 0.5, 1.6, 12.3, 1.8, [
    'Validate top drug candidates (monensin, imatinib, metformin) in macrophage infection models and animal systems to confirm host-directed therapeutic potential.',
    'Incorporate spatial transcriptomic data and infection reporters to distinguish infected from bystander cells within individual granulomas.'
], size=14, bold_prefix=False)

# Divider line
add_line(s, 0.5, 3.7, 12.3)

add_textbox(s, 0.5, 3.9, 12.3, 0.4, 'Take-Home Messages', size=16, bold=True, color=ACCENT)
add_bullet_list(s, 0.5, 4.4, 12.3, 2.8, [
    'Ensemble machine learning applied to single-cell transcriptomics can simultaneously achieve high-accuracy classification and yield biologically interpretable disease signatures.',
    'SHAP-based interpretability provides a rigorous framework for translating high-dimensional transcriptomic data into actionable therapeutic hypotheses.',
    'The computational framework is readily extensible to other infectious disease contexts where single-cell data are available.'
], size=14, bold_prefix=False)

# ============================================================
# SLIDE 15: KEY REFERENCES
# ============================================================
s = add_blank_slide()
add_header_bar(s, 'Key References')
refs = [
    '[1] WHO. Global Tuberculosis Report 2023. World Health Organization, Geneva, 2023.',
    '[2] Gideon HP, Hughes TK, Tzouanas CN, et al. Multimodal profiling of lung granulomas in macaques reveals cellular correlates of tuberculosis control. Immunity, 55(5):827-846, 2022.',
    '[3] Ke G, Meng Q, Finley T, et al. LightGBM: A highly efficient gradient boosting decision tree. Advances in Neural Information Processing Systems, 30:3146-3154, 2017.',
    '[4] Lundberg SM, Erion G, Chen H, et al. From local explanations to global understanding with explainable AI for trees. Nature Machine Intelligence, 2(1):56-67, 2020.'
]
for i, ref in enumerate(refs):
    add_textbox(s, 0.5, 1.2 + i * 1.3, 12.3, 1.2, ref, size=13, color=BODY)

# ============================================================
# SLIDE 16: THANK YOU
# ============================================================
s = add_blank_slide()
add_textbox(s, 1.5, 2.5, 10.3, 1.0, 'THANK YOU FOR LISTENING',
    size=36, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_line(s, 5.5, 3.8, 2.3)
add_textbox(s, 1.5, 4.2, 10.3, 0.5, 'Walter Odur  |  walter.odur@students.mak.ac.ug',
    size=14, color=GREY, align=PP_ALIGN.CENTER)

# ============================================================
# SAVE
# ============================================================
out_path = os.path.join(os.path.dirname(__file__), 'TB_ML_Presentation.pptx')
prs.save(out_path)
print(f'Presentation saved to: {out_path}')
print(f'Total slides: {len(prs.slides)}')
