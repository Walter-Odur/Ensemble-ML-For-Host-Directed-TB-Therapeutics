from pptx import Presentation
prs = Presentation(r'E:\MY_RESEARCH_IDEAS\ml_project\presentation\TB_ML_Presentation.pptx')
for i, slide in enumerate(prs.slides):
    texts = []
    has_img = False
    has_tbl = False
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                t = p.text.strip()
                if t:
                    texts.append(t)
        if shape.has_table:
            has_tbl = True
        try:
            _ = shape.image
            has_img = True
        except:
            pass
    first = texts[0] if texts else '(no text)'
    extras = []
    if has_img: extras.append('IMAGE')
    if has_tbl: extras.append('TABLE')
    tag = ' [' + ', '.join(extras) + ']' if extras else ''
    print(f'Slide {i+1}: {first[:90]}{tag}')
